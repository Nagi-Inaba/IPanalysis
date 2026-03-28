# -*- coding: utf-8 -*-
"""IP Analysis Studio LT Presentation — PowerPoint Generator (v3: user-spec font sizes + QR codes)"""
from __future__ import annotations

import qrcode
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Theme colors (light theme — accent teal/gold kept) ──
NAVY = RGBColor(0x0A, 0x16, 0x28)
TEAL = RGBColor(0x00, 0x96, 0x96)
TEAL_LIGHT = RGBColor(0xE6, 0xF5, 0xF5)
GOLD = RGBColor(0xC4, 0x98, 0x33)
GOLD_LIGHT = RGBColor(0xFB, 0xF5, 0xE6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0x1A, 0x1A, 0x1A)
PAGE_NUM_GRAY = RGBColor(0x88, 0x99, 0xAA)
BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_LIGHT = RGBColor(0xF4, 0xF6, 0xF8)
BORDER_LIGHT = RGBColor(0xE0, 0xE5, 0xEB)

# ── Font sizes (user spec) ──
TITLE_SIZE = 58
HEADING_SIZE = 28
BODY_SIZE = 20
TABLE_SIZE = 18
NUM_SIZE = 11  # page number — unchanged

# ── Layout constants ──
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
M = Inches(0.5)
CW = Inches(12.3)
HALF = Inches(5.9)
TOTAL_SLIDES = 9

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ── Helpers ──
def bg_white(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_WHITE


def accent_bar(slide, left, top, width=Inches(0.8), height=Pt(5), color=TEAL):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()


def slide_num(slide, num):
    tb = slide.shapes.add_textbox(SLIDE_W - Inches(1.3), Inches(0.2), Inches(1), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.text = f"{num:02d} / {TOTAL_SLIDES:02d}"
    p.font.size = Pt(NUM_SIZE)
    p.font.color.rgb = PAGE_NUM_GRAY
    p.font.name = "Noto Sans JP"
    p.alignment = PP_ALIGN.RIGHT


def txt(slide, left, top, w, h, text, size=BODY_SIZE,
        bold=False, color=TEXT, align=PP_ALIGN.LEFT, ls=1.35):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Noto Sans JP"
    p.alignment = align
    p.space_after = Pt(0)
    p.line_spacing = Pt(size * ls)
    return tb


def multiline(slide, left, top, w, h, lines, size=BODY_SIZE, color=TEXT, ls=1.4):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line.strip()
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Noto Sans JP"
        p.space_after = Pt(2)
        p.line_spacing = Pt(size * ls)
    return tb


def heading(slide, left, top, text, size=HEADING_SIZE, color=TEAL):
    return txt(slide, left, top, Inches(8), Inches(0.6), text, size=size, bold=True, color=color)


def box(slide, left, top, w, h, fill=TEAL_LIGHT, border=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    s.adjustments[0] = 0.04
    return s


def table(slide, left, top, w, rows, col_w, headers, hdr_color=TEAL):
    rh = Inches(0.45)
    ts = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, w, rh * (len(rows) + 1))
    t = ts.table
    for i, cw in enumerate(col_w):
        t.columns[i].width = Emu(int(cw * 914400))
    for i, h in enumerate(headers):
        c = t.cell(0, i)
        c.text = h
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(TABLE_SIZE)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = "Noto Sans JP"
        c.fill.solid()
        c.fill.fore_color.rgb = hdr_color
    for r, rd in enumerate(rows, 1):
        for ci, val in enumerate(rd):
            c = t.cell(r, ci)
            c.text = val
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(TABLE_SIZE)
                p.font.color.rgb = TEXT
                p.font.name = "Noto Sans JP"
                if ci == 0:
                    p.font.bold = True
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE if r % 2 == 1 else GRAY_LIGHT
    return ts


def gen_qr(url):
    qr = qrcode.QRCode(error_correction=qrcode.constants.ERROR_CORRECT_M, box_size=20, border=2)
    qr.add_data(url)
    qr.make(fit=True)
    img = qr.make_image(fill_color="black", back_color="white")
    buf = BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


# ======================================================================
# SLIDE 1: TITLE
# ======================================================================
s = prs.slides.add_slide(BLANK)
bg_white(s)

stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.15), SLIDE_H)
stripe.fill.solid()
stripe.fill.fore_color.rgb = TEAL
stripe.line.fill.background()

txt(s, Inches(0.8), Inches(1.5), Inches(9), Inches(1.5),
    "IP Analysis Studio", size=TITLE_SIZE, bold=True, color=NAVY)

txt(s, Inches(0.8), Inches(3.2), Inches(10), Inches(1.2),
    "Python スクリプトから Web アプリへ\nClaude Code で特許分析ツールを作り直した話",
    size=BODY_SIZE, color=TEXT)

txt(s, Inches(0.8), Inches(5.0), Inches(6), Inches(0.5),
    "Lightning Talk  |  知財ツール開発者オフ会", size=TABLE_SIZE, color=TEAL)

c1 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(1.5), Inches(3.5), Inches(3.5))
c1.fill.background()
c1.line.color.rgb = RGBColor(0xE0, 0xF0, 0xF0)
c1.line.width = Pt(2)
c2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.2), Inches(2.2), Inches(2.2), Inches(2.2))
c2.fill.background()
c2.line.color.rgb = RGBColor(0xE8, 0xF4, 0xF4)
c2.line.width = Pt(1.5)


# ======================================================================
# SLIDE 2: SELF-INTRO (Profile)
# ======================================================================
s = prs.slides.add_slide(BLANK)
bg_white(s)
slide_num(s, 2)

accent_bar(s, M, Inches(0.35))
txt(s, M, Inches(0.50), Inches(5), Inches(1.0),
    "自己紹介", size=TITLE_SIZE, bold=True, color=NAVY)

# Name
txt(s, M, Inches(1.5), Inches(4), Inches(0.9),
    "稲葉 凪", size=TITLE_SIZE, bold=True, color=NAVY)
txt(s, Inches(4.5), Inches(1.7), Inches(3), Inches(0.5),
    "Nagi Inaba", size=BODY_SIZE, color=TEXT)

# Left column
col_l = M
y = Inches(2.5)
heading(s, col_l, y, "所属")
y += Inches(0.5)
multiline(s, col_l, y, Inches(5.5), Inches(0.8), [
    "大阪工業大学大学院 知的財産研究科",
    "小林誠ゼミ（シクロハイジア）",
])
y += Inches(0.85)
heading(s, col_l, y, "経歴")
y += Inches(0.5)
multiline(s, col_l, y, Inches(5.5), Inches(0.8), [
    "沼津高専 電気電子工学科 卒業",
    "大阪工業大学 知的財産学部 卒業",
])
y += Inches(0.85)
heading(s, col_l, y, "専門")
y += Inches(0.5)
multiline(s, col_l, y, Inches(5.5), Inches(1.2), [
    "特許情報分析（IPランドスケープ）",
    "知財戦略・M&A知財評価",
    "ツール開発（Python / Web アプリ(Next.js)）",
])

# Right column
col_r = Inches(6.5)
yr = Inches(2.5)
heading(s, col_r, yr, "資格")
yr += Inches(0.5)
multiline(s, col_r, yr, Inches(6.3), Inches(1.2), [
    "2級知的財産管理技能士",
    "知的財産アナリスト（特許）— 3月31日 合格発表予定",
    "紅茶検定上級",
])
yr += Inches(1.3)
heading(s, col_r, yr, "実務経験")
yr += Inches(0.5)
multiline(s, col_r, yr, Inches(6.3), Inches(1.0), [
    "弁理士事務所インターン（業務効率化ツール開発・特許出願）",
    "知財コンサルティングファーム インターン（IPランドスケープ実践）",
])


# ======================================================================
# SLIDE 3: SELF-INTRO (Projects)
# ======================================================================
s = prs.slides.add_slide(BLANK)
bg_white(s)
slide_num(s, 3)

accent_bar(s, M, Inches(0.35))
txt(s, M, Inches(0.50), Inches(5), Inches(1.0),
    "個人開発", size=TITLE_SIZE, bold=True, color=NAVY)

projects = [
    ("万博予約自動化", "大阪万博来場予約自動化スクリプトの作成・販売", GOLD_LIGHT, GOLD),
    ("街頭演説マップ", "選挙の街頭演説を地図上で可視化。誰がいつ実施中・予定なのかをピンポイントで表示（チームみらい用）", TEAL_LIGHT, TEAL),
    ("ポスティングマップ", "国勢調査のデータを使用し、ターゲットが多いエリアを探して効率よくポスティングできる Web アプリ", TEAL_LIGHT, TEAL),
]

y = Inches(1.8)
for title, desc, bg, accent in projects:
    box(s, M, y, CW, Inches(1.4), fill=bg)
    txt(s, M + Inches(0.3), y + Inches(0.15), Inches(11), Inches(0.5),
        title, size=HEADING_SIZE, bold=True, color=accent)
    txt(s, M + Inches(0.3), y + Inches(0.7), Inches(11), Inches(0.5),
        desc, size=BODY_SIZE, color=TEXT)
    y += Inches(1.6)


# ======================================================================
# SLIDE 4: BEFORE
# ======================================================================
s = prs.slides.add_slide(BLANK)
bg_white(s)
slide_num(s, 4)

accent_bar(s, M, Inches(0.35), color=GOLD)
txt(s, M, Inches(0.50), Inches(3), Inches(1.0),
    "Before", size=TITLE_SIZE, bold=True, color=NAVY)
txt(s, Inches(4.2), Inches(0.70), Inches(6), Inches(0.5),
    "Python スクリプト時代", size=BODY_SIZE, color=GOLD)

# Left column
col_l = M
y = Inches(1.7)
heading(s, col_l, y, "やっていたこと")
y += Inches(0.5)
multiline(s, col_l, y, Inches(5.8), Inches(1.5), [
    "CKSweb から CSV をダウンロード",
    "Python スクリプトで CSV を整理・集計",
    "特許分類増減率の計算、出願人カウントなどを pandas で処理",
])
y += Inches(1.6)
heading(s, col_l, y, "毎回手作業だったこと", color=GOLD)
y += Inches(0.5)
multiline(s, col_l, y, Inches(5.8), Inches(1.5), [
    "Excel でテーブルを範囲選択してグラフ化",
    "グラフの書式設定（軸ラベル、色、凡例）を毎回手動",
    "分析のたびに同じ作業を繰り返す",
])

# Right: Issues
col_r = Inches(7.0)
yr = Inches(1.7)
heading(s, col_r, yr, "課題", color=GOLD)
yr += Inches(0.55)

issues = [
    "集計ロジックは自動化できたが、\n可視化が毎回手作業",
    "分析結果を他人に共有するのが面倒\n（Excel ファイルを渡すしかない）",
    "コードが 1ファイル 1,000行超 に肥大化",
]
for issue in issues:
    box(s, col_r, yr, Inches(5.8), Inches(1.2), fill=GOLD_LIGHT, border=GOLD)
    txt(s, col_r + Inches(0.3), yr + Inches(0.15), Inches(5.2), Inches(0.9),
        issue, size=BODY_SIZE, color=TEXT)
    yr += Inches(1.4)


# ======================================================================
# SLIDE 5: AFTER
# ======================================================================
s = prs.slides.add_slide(BLANK)
bg_white(s)
slide_num(s, 5)

accent_bar(s, M, Inches(0.35))
txt(s, M, Inches(0.50), Inches(3), Inches(1.0),
    "After", size=TITLE_SIZE, bold=True, color=NAVY)
txt(s, Inches(3.5), Inches(0.70), Inches(8), Inches(0.5),
    "Claude Code で Web アプリ化", size=BODY_SIZE, color=TEAL)

items = [
    ("1", "Streamlit で Web UI 化",
     "ファイルアップロード → 前処理 → 集計 → グラフ の3ステップ。ブラウザで完結"),
    ("2", "Altair でグラフ自動生成",
     "15種以上のインタラクティブチャート。ツールチップ、ズーム、フィルタが標準装備"),
    ("3", "データ形式の自動判定",
     "CKSweb / J-PlatPat / Questel の列名パターンで自動検出"),
    ("4", "出願人名寄せの内蔵",
     "JSON 辞書ベースで表記揺れを統一。GUI で編集・インポート・エクスポート可能"),
    ("5", "1,069行 → 4ファイル分割リファクタリング",
     "app.py (455行) + charts.py (555行) + cached_agg.py (70行) + constants.py (44行)"),
]

y = Inches(1.7)
for num, title, desc in items:
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, M, y + Inches(0.05), Inches(0.5), Inches(0.5))
    circ.fill.solid()
    circ.fill.fore_color.rgb = TEAL
    circ.line.fill.background()
    p = circ.text_frame.paragraphs[0]
    p.text = num
    p.font.size = Pt(BODY_SIZE)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Noto Sans JP"
    p.alignment = PP_ALIGN.CENTER

    txt(s, Inches(1.2), y, Inches(11), Inches(0.45),
        title, size=BODY_SIZE, bold=True, color=NAVY)
    txt(s, Inches(1.2), y + Inches(0.45), Inches(11), Inches(0.4),
        desc, size=TABLE_SIZE, color=TEXT)
    y += Inches(1.05)


# ======================================================================
# SLIDE 6: TECH STACK
# ======================================================================
s = prs.slides.add_slide(BLANK)
bg_white(s)
slide_num(s, 6)

accent_bar(s, M, Inches(0.35))
txt(s, M, Inches(0.50), Inches(5), Inches(1.0),
    "技術スタック", size=TITLE_SIZE, bold=True, color=NAVY)

# Core table
heading(s, M, Inches(1.5), "コア")
table(s, M, Inches(2.0), Inches(8), [
    ["Python 3.8+", "言語", "特許分析ライブラリとの親和性"],
    ["Streamlit", "Web UI", "Python だけで UI が書ける"],
    ["pandas", "データ処理", "CSV/Excel の読み込み・集計"],
    ["openpyxl", "Excel I/O", "xlsx の読み書き"],
    ["Altair", "可視化", "宣言的 API、Vega-Lite ベース"],
], [2.2, 1.8, 4.0], ["技術", "用途", "選定理由"])

# Infra
heading(s, M, Inches(5.0), "インフラ・運用")
table(s, M, Inches(5.5), Inches(8), [
    ["Streamlit Cloud", "ホスティング（GitHub 連携自動デプロイ）"],
    ["pytest", "テスト（31テスト）"],
    ["GitHub", "ソースコード管理"],
], [2.2, 5.8], ["技術", "用途"])

# Dev tools
heading(s, Inches(9.2), Inches(1.5), "開発ツール")
table(s, Inches(9.2), Inches(2.0), Inches(3.6), [
    ["Claude Code", "コーディング・リファクタ・テスト生成"],
    ["Git", "バージョン管理"],
], [1.5, 2.1], ["ツール", "役割"])


# ======================================================================
# SLIDE 7: 17 ANALYSIS FEATURES
# ======================================================================
s = prs.slides.add_slide(BLANK)
bg_white(s)
slide_num(s, 7)

accent_bar(s, M, Inches(0.35))
txt(s, M, Inches(0.50), Inches(4), Inches(1.0),
    "分析機能", size=TITLE_SIZE, bold=True, color=NAVY)
txt(s, Inches(4.2), Inches(0.40), Inches(1.5), Inches(1.0),
    "17", size=48, bold=True, color=TEAL)
txt(s, Inches(5.5), Inches(0.65), Inches(1), Inches(0.5),
    "種", size=HEADING_SIZE, bold=True, color=NAVY)

categories = [
    ("出願動向", [
        "出願件数推移",
    ]),
    ("特許分類 / Fターム", [
        "特許分類増減率（バブル）",
        "特許分類集計 / メイングループ",
        "分類別 年次ヒートマップ",
        "分類ツリーマップ",
        "Fターム分布",
        "Fターム年次ヒートマップ",
    ]),
    ("出願人", [
        "ランキング（筆頭 / 全体）",
        "増減率 / 年次推移 / シェア",
        "参入退出分析",
        "出願人 × 分類 ヒートマップ",
        "共同出願ネットワーク",
    ]),
    ("引用", [
        "被引用ポジショニングマップ",
        "被引用出願一覧",
    ]),
]

x_pos = [M, Inches(3.5), Inches(7.0), Inches(10.5)]
col_w = Inches(2.8)
for i, (cat, items_list) in enumerate(categories):
    x = x_pos[i]
    y = Inches(1.6)
    box(s, x, y, col_w, Inches(0.55), fill=TEAL_LIGHT)
    txt(s, x + Inches(0.15), y + Inches(0.07), col_w - Inches(0.3), Inches(0.4),
        cat, size=TABLE_SIZE, bold=True, color=TEAL)
    y += Inches(0.7)
    for item in items_list:
        txt(s, x + Inches(0.15), y, col_w - Inches(0.2), Inches(0.4),
            "・ " + item, size=TABLE_SIZE, color=TEXT)
        y += Inches(0.45)


# ======================================================================
# SLIDE 8: LEARNINGS
# ======================================================================
s = prs.slides.add_slide(BLANK)
bg_white(s)
slide_num(s, 8)

accent_bar(s, M, Inches(0.35))
txt(s, M, Inches(0.50), Inches(12), Inches(1.0),
    "Claude Code で開発して感じたこと", size=TITLE_SIZE, bold=True, color=NAVY)

# Left: Good
col_l = M
y = Inches(1.7)
heading(s, col_l, y, "よかったこと")
y += Inches(0.55)

good = [
    ("プロトタイプが爆速", "Streamlit + Altair の組み合わせを提案してくれた"),
    ("リファクタリングが楽", "1,069行 → 4ファイル分割を一発で実行"),
    ("テスト生成", "分析ロジックの pytest を自動生成（31テスト）"),
    ("データ形式対応の追加", "Questel Orbit 対応を列名パターンから自動実装"),
]
for t, d in good:
    box(s, col_l, y, Inches(5.8), Inches(1.05), fill=TEAL_LIGHT)
    txt(s, col_l + Inches(0.25), y + Inches(0.1), Inches(5.3), Inches(0.45),
        t, size=BODY_SIZE, bold=True, color=NAVY)
    txt(s, col_l + Inches(0.25), y + Inches(0.55), Inches(5.3), Inches(0.4),
        d, size=TABLE_SIZE, color=TEXT)
    y += Inches(1.15)

# Right: Cautions
col_r = Inches(6.8)
yr = Inches(1.7)
heading(s, col_r, yr, "気をつけていること", color=GOLD)
yr += Inches(0.55)

caution = [
    ("動作確認は必ず自分で", "生成コードは必ず動作確認してから採用"),
    ("ドメイン知識は人間が持つ", "特許分析の文脈は Claude が知らない領域"),
    ("計画 → テスト → 実装", "大きな変更は段階的に進める"),
]
for t, d in caution:
    box(s, col_r, yr, Inches(5.8), Inches(1.05), fill=GOLD_LIGHT)
    txt(s, col_r + Inches(0.25), yr + Inches(0.1), Inches(5.3), Inches(0.45),
        t, size=BODY_SIZE, bold=True, color=NAVY)
    txt(s, col_r + Inches(0.25), yr + Inches(0.55), Inches(5.3), Inches(0.4),
        d, size=TABLE_SIZE, color=TEXT)
    yr += Inches(1.15)


# ======================================================================
# SLIDE 9: DEMO & REPO with QR CODES
# ======================================================================
s = prs.slides.add_slide(BLANK)
bg_white(s)
slide_num(s, 9)

accent_bar(s, Inches(5.6), Inches(0.8), width=Inches(1.2))
txt(s, Inches(2), Inches(1.0), Inches(9), Inches(1.0),
    "Demo & Repository", size=TITLE_SIZE, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# --- QR codes ---
qr_size = Inches(2.5)
qr_y = Inches(2.3)

# Demo QR (left)
demo_qr = gen_qr("https://ipanalysis-webapp.streamlit.app/")
s.shapes.add_picture(demo_qr, Inches(2.2), qr_y, qr_size, qr_size)
txt(s, Inches(0.5), qr_y - Inches(0.5), Inches(5.5), Inches(0.5),
    "ライブデモ", size=HEADING_SIZE, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
txt(s, Inches(0.5), qr_y + qr_size + Inches(0.2), Inches(5.5), Inches(0.5),
    "https://ipanalysis-webapp.streamlit.app/", size=TABLE_SIZE, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
txt(s, Inches(0.5), qr_y + qr_size + Inches(0.65), Inches(5.5), Inches(0.4),
    "アカウント不要・無料・サンプルデータ内蔵", size=TABLE_SIZE, color=TEXT, align=PP_ALIGN.CENTER)

# GitHub QR (right)
gh_qr = gen_qr("https://github.com/Nagi-Inaba/IPanalysis")
s.shapes.add_picture(gh_qr, Inches(8.5), qr_y, qr_size, qr_size)
txt(s, Inches(6.8), qr_y - Inches(0.5), Inches(5.5), Inches(0.5),
    "GitHub", size=HEADING_SIZE, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
txt(s, Inches(6.8), qr_y + qr_size + Inches(0.2), Inches(5.5), Inches(0.5),
    "https://github.com/Nagi-Inaba/IPanalysis", size=TABLE_SIZE, bold=True, color=TEXT, align=PP_ALIGN.CENTER)

# Footer
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(6.3), Inches(2), Pt(1))
line.fill.solid()
line.fill.fore_color.rgb = BORDER_LIGHT
line.line.fill.background()

txt(s, Inches(2), Inches(6.5), Inches(9), Inches(0.5),
    "稲葉 凪 — 大阪工業大学大学院 知的財産研究科", size=TABLE_SIZE, color=TEXT, align=PP_ALIGN.CENTER)


# ======================================================================
# SAVE
# ======================================================================
out = r"C:\Users\nagii\OneDrive\ドキュメント\GitHub\IPanalysis\slides\lt-presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
